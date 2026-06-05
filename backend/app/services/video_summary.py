from __future__ import annotations

from dataclasses import asdict, dataclass, replace
from pathlib import Path
import json
import mimetypes
import os
import re
import subprocess
from textwrap import wrap
from urllib.error import HTTPError, URLError
from urllib.parse import quote
from urllib.request import Request, urlopen


DEFAULT_STORAGE_BUCKET = "video-summaries"

# Don't delete following commands
# [GenAI Usage - concurrency]: Codex Prompt
# Right now, Backend video summary generation has 2 problems:
# 1. If two users generate a summary video at the same time, they will overwrite each other because right now the folder is stored by paper-id only.
# 2. Even though video being generated already, the video summary is not read first from Supabase bucket in the fetch logic.I the same person needs to re-run video generation again after refreshing the page.
# In all, since the video summary is the same for all user, I would love to keep the same storage structure. All users shall the same video summary, because there is  no personalization here. So, add locking mechanisms to provide two users concurrentlly generate same video summray. Second, before every generation, load from supabase bucket to see if there is one already
# [GenAI Usage]: Codex Prompt
# After analyze with the current pattern, let's don't worry about generation caching right now. it is too complicated and out-of-scope. It's a beta-feature that would be work on later. i think a GET method would be fine
# [GenAI Usage]: Codex Response Begins
DEFAULT_STORAGE_PREFIX = "video_summaries"
VIDEO_SUMMARY_MANIFEST_FILENAME = "video_summary_manifest.json"
VIDEO_SUMMARY_MANIFEST_VERSION = 1
DEFAULT_SCRATCH_ROOT = Path(os.getenv("VIDEO_SUMMARY_SCRATCH_ROOT", "/tmp/generated"))
DEFAULT_OUTPUT_ROOT = DEFAULT_SCRATCH_ROOT / "video_summaries"
SLIDE_WIDTH = 1280
SLIDE_HEIGHT = 720
VIDEO_FPS = 2
MAX_VIDEO_SLIDES = 8


#[GenAI Usage] Prompt: implement the video summary functionality. Make it similar to text summary, but instead use a python package that directly creates presentation slides, . Write the script with timestamps. Then play the slides with the timestamps and render a video.
#[GenAI Usage] LLM response begins

#[GenAI Usage] Codex prompt: Why frontend got that error when i click on generate video summary.
# [Detailed error info: OSError: [Errno 30] Read-only file system: 'generated' on Vercel].
# How to fix the current video generation pipeline on vercel deployment server.
#[GenAI Usage] Codex response begins:

class VideoSummaryError(Exception):
    pass


@dataclass(frozen=True)
class VideoSlide:
    title: str
    bullets: list[str]
    narration: str
    subtitle: str
    duration_seconds: int
    visual_asset_index: int | None = None
    visual_path: str | None = None
    visual_caption: str | None = None
    visual_reason: str | None = None


@dataclass(frozen=True)
class TimestampedScriptLine:
    start_seconds: int
    end_seconds: int
    slide_title: str
    narration: str
    subtitle: str


@dataclass(frozen=True)
class VideoSummaryArtifact:
    output_dir: str
    pptx_path: str
    video_path: str
    silent_video_path: str | None
    audio_path: str | None
    script_path: str
    script_json_path: str
    subtitle_vtt_path: str
    subtitle_srt_path: str
    asset_metadata_path: str
    slide_image_paths: list[str]
    visual_asset_paths: list[str]
    slides: list[dict]
    script: list[dict]
    voiceover: dict | None

    def to_dict(self) -> dict:
        return asdict(self)


def load_video_summary_from_storage(paper_id: str) -> VideoSummaryArtifact | None:
    supabase_url, service_key, bucket = _storage_config()
    object_path = (
        f"{DEFAULT_STORAGE_PREFIX}/{paper_id}/{VIDEO_SUMMARY_MANIFEST_FILENAME}"
    )
    payload = _download_storage_json(
        supabase_url=supabase_url,
        service_key=service_key,
        bucket=bucket,
        object_path=object_path,
    )
    if payload is None:
        return None
    if payload.get("manifest_version") != VIDEO_SUMMARY_MANIFEST_VERSION:
        raise VideoSummaryError("Stored video summary manifest has an unsupported version.")

    artifact = payload.get("video_summary")
    if not isinstance(artifact, dict):
        raise VideoSummaryError("Stored video summary manifest is invalid.")
    try:
        return VideoSummaryArtifact(**artifact)
    except TypeError as exc:
        raise VideoSummaryError("Stored video summary manifest is invalid.") from exc


def create_video_summary_artifacts(
    *,
    paper_id: str,
    title: str,
    summary_text: str,
    slide_plan: list[dict] | None = None,
    image_assets: list | None = None,
    video_instructions: str | None = None,
    slide_duration_seconds: int = 8,
    include_voiceover: bool = False,
    output_root: Path = DEFAULT_OUTPUT_ROOT,
) -> VideoSummaryArtifact:
    if not summary_text or not summary_text.strip():
        raise VideoSummaryError("Summary text is required before making a video.")

    output_dir = output_root / paper_id
    output_dir.mkdir(parents=True, exist_ok=True)

    visual_assets = image_assets or []
    slides = _slides_from_plan(
        title=title,
        slide_plan=slide_plan,
        image_assets=visual_assets,
        summary_text=summary_text,
        video_instructions=video_instructions,
        slide_duration_seconds=slide_duration_seconds,
    )
    voiceover = None
    if include_voiceover:
        try:
            from app.services.voiceover import (
                VoiceoverError,
                generate_elevenlabs_voiceover,
                voiceover_slide_durations,
            )
            voiceover = generate_elevenlabs_voiceover(
                slides,
                output_dir / "audio",
            )
            slides = _apply_voiceover_durations(
                slides,
                voiceover_slide_durations(voiceover),
            )
        except VoiceoverError as exc:
            raise VideoSummaryError(str(exc)) from exc

    script = _build_timestamped_script(slides)

    pptx_path = output_dir / "video_summary.pptx"
    script_path = output_dir / "video_summary_script.txt"
    script_json_path = output_dir / "video_summary_script.json"
    subtitle_vtt_path = output_dir / "video_summary_subtitles.vtt"
    subtitle_srt_path = output_dir / "video_summary_subtitles.srt"
    asset_metadata_path = output_dir / "video_summary_assets.json"
    video_path = output_dir / "video_summary.mp4"
    silent_video_path = output_dir / "video_summary_silent.mp4" if voiceover else None
    frames_dir = output_dir / "frames"
    frames_dir.mkdir(exist_ok=True)

    _write_presentation(slides, pptx_path)
    slide_image_paths = _render_slide_images(slides, frames_dir)
    _write_script(script, script_path)
    _write_script_json(script, script_json_path)
    _write_vtt(script, subtitle_vtt_path)
    _write_srt(script, subtitle_srt_path)
    _write_asset_metadata(visual_assets, asset_metadata_path)
    rendered_video_path = silent_video_path or video_path
    _render_video(slide_image_paths, slides, rendered_video_path)
    if voiceover:
        _mux_video_audio(
            rendered_video_path,
            Path(voiceover.combined_audio_path),
            video_path,
        )

    static_root = output_dir
    storage_object_prefix = f"{DEFAULT_STORAGE_PREFIX}/{paper_id}"
    storage_urls = _upload_artifacts_to_storage(
        _collect_artifact_paths(
            pptx_path=pptx_path,
            video_path=video_path,
            silent_video_path=silent_video_path,
            voiceover=voiceover,
            script_path=script_path,
            script_json_path=script_json_path,
            subtitle_vtt_path=subtitle_vtt_path,
            subtitle_srt_path=subtitle_srt_path,
            asset_metadata_path=asset_metadata_path,
            slide_image_paths=slide_image_paths,
            visual_assets=visual_assets,
        ),
        static_root=static_root,
        object_prefix=storage_object_prefix,
    )

    video_url = _artifact_url(video_path, static_root, storage_urls)
    # parental path
    out_dir = video_url.rsplit("/", 1)[0] + "/"
    artifact = VideoSummaryArtifact(
        output_dir=out_dir,
        pptx_path=_artifact_url(pptx_path, static_root, storage_urls),
        video_path=video_url,
        silent_video_path=(
            _artifact_url(silent_video_path, static_root, storage_urls)
            if silent_video_path
            else None
        ),
        audio_path=(
            _artifact_url(Path(voiceover.combined_audio_path), static_root, storage_urls)
            if voiceover
            else None
        ),
        script_path=_artifact_url(script_path, static_root, storage_urls),
        script_json_path=_artifact_url(script_json_path, static_root, storage_urls),
        subtitle_vtt_path=_artifact_url(subtitle_vtt_path, static_root, storage_urls),
        subtitle_srt_path=_artifact_url(subtitle_srt_path, static_root, storage_urls),
        asset_metadata_path=_artifact_url(asset_metadata_path, static_root, storage_urls),
        slide_image_paths=[
            _artifact_url(path, static_root, storage_urls)
            for path in slide_image_paths
        ],
        visual_asset_paths=[
            _artifact_url(Path(asset.path), static_root, storage_urls)
            for asset in visual_assets
        ],
        slides=[
            _public_slide_dict(slide, static_root, storage_urls)
            for slide in slides
        ],
        script=[asdict(line) for line in script],
        voiceover=_public_voiceover_dict(voiceover, static_root, storage_urls) if voiceover else None,
    )
    _upload_video_summary_manifest(
        paper_id=paper_id,
        artifact=artifact,
        output_dir=output_dir,
        static_root=static_root,
        object_prefix=storage_object_prefix,
    )
    return artifact


def _artifact_url(path: Path, static_root: Path, storage_urls: dict[Path, str]) -> str:
    storage_url = storage_urls.get(path)
    if storage_url:
        return storage_url
    raise VideoSummaryError(f"Generated artifact was not uploaded to durable storage: {path.name}")


def _public_slide_dict(
    slide: VideoSlide,
    static_root: Path,
    storage_urls: dict[Path, str],
) -> dict:
    data = asdict(slide)
    if slide.visual_path:
        data["visual_path"] = _artifact_url(Path(slide.visual_path), static_root, storage_urls)
    return data


def _public_voiceover_dict(voiceover, static_root: Path, storage_urls: dict[Path, str]) -> dict:
    data = voiceover.to_dict()
    data["combined_audio_path"] = _artifact_url(
        Path(voiceover.combined_audio_path),
        static_root,
        storage_urls,
    )
    for segment in data.get("segments", []):
        if segment.get("audio_path"):
            segment["audio_path"] = _artifact_url(
                Path(segment["audio_path"]),
                static_root,
                storage_urls,
            )
    return data


def _collect_artifact_paths(
    *,
    pptx_path: Path,
    video_path: Path,
    silent_video_path: Path | None,
    voiceover,
    script_path: Path,
    script_json_path: Path,
    subtitle_vtt_path: Path,
    subtitle_srt_path: Path,
    asset_metadata_path: Path,
    slide_image_paths: list[Path],
    visual_assets: list,
) -> list[Path]:
    paths = [
        pptx_path,
        video_path,
        script_path,
        script_json_path,
        subtitle_vtt_path,
        subtitle_srt_path,
        asset_metadata_path,
        *slide_image_paths,
        *[Path(asset.path) for asset in visual_assets],
    ]
    if silent_video_path:
        paths.append(silent_video_path)
    if voiceover:
        paths.append(Path(voiceover.combined_audio_path))
        paths.extend(Path(segment.audio_path) for segment in voiceover.segments)

    seen = set()
    unique_paths = []
    for path in paths:
        if path in seen or not path.exists() or not path.is_file():
            continue
        seen.add(path)
        unique_paths.append(path)
    return unique_paths


def _upload_artifacts_to_storage(
    paths: list[Path],
    *,
    static_root: Path,
    object_prefix: str,
) -> dict[Path, str]:
    supabase_url, service_key, bucket = _storage_config()

    uploaded_urls: dict[Path, str] = {}
    for path in paths:
        object_path = _storage_object_path(path, static_root, object_prefix)
        _upload_file_to_storage(
            supabase_url=supabase_url,
            service_key=service_key,
            bucket=bucket,
            object_path=object_path,
            file_path=path,
        )
        uploaded_urls[path] = _storage_public_url(supabase_url, bucket, object_path)
    return uploaded_urls


def _upload_video_summary_manifest(
    *,
    paper_id: str,
    artifact: VideoSummaryArtifact,
    output_dir: Path,
    static_root: Path,
    object_prefix: str,
) -> None:
    manifest_path = output_dir / VIDEO_SUMMARY_MANIFEST_FILENAME
    manifest_path.write_text(
        json.dumps(
            {
                "manifest_version": VIDEO_SUMMARY_MANIFEST_VERSION,
                "paper_id": paper_id,
                "video_summary": artifact.to_dict(),
            },
            indent=2,
        ),
        encoding="utf-8",
    )
    _upload_artifacts_to_storage(
        [manifest_path],
        static_root=static_root,
        object_prefix=object_prefix,
    )


def _storage_config() -> tuple[str, str, str]:
    supabase_url = (os.getenv("SUPABASE_URL") or "").rstrip("/")
    service_key = os.getenv("SUPABASE_SERVICE_ROLE_KEY")
    bucket = os.getenv("VIDEO_SUMMARY_STORAGE_BUCKET", DEFAULT_STORAGE_BUCKET)
    if not supabase_url or not service_key:
        raise VideoSummaryError(
            "SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY are required to persist video summaries."
        )
    return supabase_url, service_key, bucket


def _storage_object_path(path: Path, static_root: Path, object_prefix: str) -> str:
    try:
        relative_path = path.relative_to(static_root).as_posix()
    except ValueError:
        relative_path = path.name
    return f"{object_prefix.rstrip('/')}/{relative_path}"


def _upload_file_to_storage(
    *,
    supabase_url: str,
    service_key: str,
    bucket: str,
    object_path: str,
    file_path: Path,
) -> None:
    content_type = _content_type_for_file(file_path)
    url = (
        f"{supabase_url}/storage/v1/object/"
        f"{quote(bucket, safe='')}/{quote(object_path, safe='/')}"
    )
    request = Request(
        url,
        data=file_path.read_bytes(),
        method="POST",
        headers={
            "Authorization": f"Bearer {service_key}",
            "apikey": service_key,
            "Content-Type": content_type,
            "Cache-Control": "3600",
            "x-upsert": "true",
        },
    )
    try:
        with urlopen(request, timeout=60) as response:
            response.read()
    except HTTPError as exc:
        body = exc.read(500).decode("utf-8", errors="replace")
        raise VideoSummaryError(
            f"Supabase Storage upload failed for {object_path} with HTTP {exc.code}: {body}"
        ) from exc
    except (TimeoutError, URLError, OSError) as exc:
        raise VideoSummaryError(
            f"Supabase Storage upload failed for {object_path}."
        ) from exc


def _download_storage_json(
    *,
    supabase_url: str,
    service_key: str,
    bucket: str,
    object_path: str,
) -> dict | None:
    url = (
        f"{supabase_url}/storage/v1/object/"
        f"{quote(bucket, safe='')}/{quote(object_path, safe='/')}"
    )
    request = Request(
        url,
        method="GET",
        headers={
            "Authorization": f"Bearer {service_key}",
            "apikey": service_key,
        },
    )
    try:
        with urlopen(request, timeout=30) as response:
            payload = json.loads(response.read().decode("utf-8"))
    except HTTPError as exc:
        body = exc.read(500).decode("utf-8", errors="replace")
        if exc.code == 404 or (
            exc.code == 400 and "not found" in body.lower()
        ):
            return None
        raise VideoSummaryError(
            f"Supabase Storage read failed for {object_path} with HTTP {exc.code}: {body}"
        ) from exc
    except (json.JSONDecodeError, UnicodeDecodeError) as exc:
        raise VideoSummaryError("Stored video summary manifest is invalid.") from exc
    except (TimeoutError, URLError, OSError) as exc:
        raise VideoSummaryError(
            f"Supabase Storage read failed for {object_path}."
        ) from exc

    if not isinstance(payload, dict):
        raise VideoSummaryError("Stored video summary manifest is invalid.")
    return payload


# [GenAI Usage - concurrency]: Codex Response Ends
# [GenAI Usage - concurrency]: Reflection
# After inspecting the result about "official solution to fix concurrency" from Codex, I decided not to add generation locking mechanism because it adds unnecessary complexity --  we can tolerant the inconsistency right now because collision secenaios are rare and unlikely to happen. 
# There is no need to pay such high design cost to resolve small bugs that are unlikely to happen.
# One safeguard we also introduce is that every user would have its own personal namespace during temporary directory to avoid overwrite. 
# Then, the backend now provides GET /papers/{paper_id}/video-summary, which reads the completed video summary manifest from Supabase Storage. POST method remains intact and unchanged -- persisting storage in Supabase bucket. 
# Therefore, the code shall be accepted right now.

def _storage_public_url(supabase_url: str, bucket: str, object_path: str) -> str:
    return (
        f"{supabase_url}/storage/v1/object/public/"
        f"{quote(bucket, safe='')}/{quote(object_path, safe='/')}"
    )


def _content_type_for_file(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".vtt":
        return "text/vtt"
    if suffix == ".srt":
        return "application/x-subrip"
    return mimetypes.guess_type(file_path.name)[0] or "application/octet-stream"

# Don't delete or modify this comment
# [GenAI Usage] Reflection:
# When I test this feature on deployed Vercel server, I got an internal server error 500, because the deployed FastAPI route crashed due to OS exception. 
# The cause is that the legacy generation pipeline need to create new files in ./generated/video_summaries folder, which is read-only.
# The fix is to treat local files as scratch artifacts only: render under /tmp on Vercel, upload the final video/PPTX/scripts/subtitles/images/audio to durable Supabase file system.
# Alongside changes in this file, Codex will create a robust migration sql to define and create a new storage bucket to persist video summary.
# Finally, return the storage url to frontend for rendering
# The code shall be accepted right now, although there might be pending consistency issue. However, it suffices to be a good starting point. We will fix the upcoming issue later when encountered. We would imagine and endure some failures in production, because it is a beta-feature.
def _slides_from_plan(
    *,
    title: str,
    slide_plan: list[dict] | None,
    image_assets: list,
    summary_text: str,
    video_instructions: str | None,
    slide_duration_seconds: int,
) -> list[VideoSlide]:
    if not slide_plan:
        return _build_fallback_slides(
            title=title,
            summary_text=summary_text,
            video_instructions=video_instructions,
            slide_duration_seconds=slide_duration_seconds,
            image_assets=image_assets,
        )

    asset_by_index = {asset.index: asset for asset in image_assets}
    slides: list[VideoSlide] = []
    default_duration = _clamp_duration(slide_duration_seconds)

    for index, raw_slide in enumerate(slide_plan[:MAX_VIDEO_SLIDES]):
        title_text = _clean_text(raw_slide.get("title"), limit=72) or f"Slide {index + 1}"
        bullets = _clean_bullets(raw_slide.get("bullets"))
        narration = _clean_text(raw_slide.get("narration"), limit=900)
        if not narration:
            narration = _bullets_to_narration(title_text, bullets, video_instructions)
        subtitle = _clean_text(raw_slide.get("subtitle"), limit=260) or _short_caption(narration)
        duration = _coerce_duration(raw_slide.get("duration_seconds"), default_duration)

        visual_asset_index = _coerce_int(raw_slide.get("visual_asset_index"))
        visual_asset = asset_by_index.get(visual_asset_index)
        if visual_asset is None and image_assets:
            visual_asset = image_assets[index % len(image_assets)]
            visual_asset_index = visual_asset.index

        slides.append(
            VideoSlide(
                title=title_text,
                bullets=bullets or [_short_caption(narration, limit=110)],
                narration=narration,
                subtitle=subtitle,
                duration_seconds=duration,
                visual_asset_index=visual_asset_index if visual_asset else None,
                visual_path=visual_asset.path if visual_asset else None,
                visual_caption=_clean_text(
                    raw_slide.get("visual_caption")
                    or (visual_asset.label if visual_asset else None),
                    limit=150,
                ),
                visual_reason=_clean_text(raw_slide.get("visual_reason"), limit=220),
            )
        )

    return slides or _build_fallback_slides(
        title=title,
        summary_text=summary_text,
        video_instructions=video_instructions,
        slide_duration_seconds=slide_duration_seconds,
        image_assets=image_assets,
    )


def _build_fallback_slides(
    *,
    title: str,
    summary_text: str,
    video_instructions: str | None,
    slide_duration_seconds: int,
    image_assets: list,
) -> list[VideoSlide]:
    duration = _clamp_duration(slide_duration_seconds)
    sections = _summary_to_sections(summary_text)
    slides: list[VideoSlide] = [
        VideoSlide(
            title="Paper at a Glance",
            bullets=[_clean_text(title, limit=120)],
            narration=f"This video summarizes the paper: {title}.",
            subtitle=f"Summary of {title}",
            duration_seconds=duration,
            visual_asset_index=image_assets[0].index if image_assets else None,
            visual_path=image_assets[0].path if image_assets else None,
            visual_caption=image_assets[0].label if image_assets else None,
        )
    ]

    for index, (heading, content) in enumerate(sections[:5], start=1):
        bullets = _content_to_bullets(content)
        if not bullets:
            continue
        visual_asset = image_assets[index % len(image_assets)] if image_assets else None
        narration = _bullets_to_narration(heading, bullets, video_instructions)
        slides.append(
            VideoSlide(
                title=heading,
                bullets=bullets,
                narration=narration,
                subtitle=_short_caption(narration),
                duration_seconds=_estimate_duration(narration, duration),
                visual_asset_index=visual_asset.index if visual_asset else None,
                visual_path=visual_asset.path if visual_asset else None,
                visual_caption=visual_asset.label if visual_asset else None,
            )
        )

    return slides


def _summary_to_sections(summary_text: str) -> list[tuple[str, str]]:
    sections: list[tuple[str, str]] = []
    current_heading = "Summary"
    current_lines: list[str] = []

    for line in summary_text.splitlines():
        clean_line = line.strip()
        heading_match = re.match(r"^#{1,4}\s+(.+)$", clean_line)
        if heading_match:
            if current_lines:
                sections.append((current_heading, "\n".join(current_lines).strip()))
                current_lines = []
            current_heading = heading_match.group(1).strip()
            continue
        current_lines.append(clean_line)

    if current_lines:
        sections.append((current_heading, "\n".join(current_lines).strip()))

    return sections or [("Summary", summary_text)]


def _content_to_bullets(content: str, *, max_bullets: int = 3) -> list[str]:
    bullets = []
    for line in content.splitlines():
        clean_line = re.sub(r"^[\-*•]\s*", "", line).strip()
        clean_line = re.sub(r"^\d+\.\s*", "", clean_line)
        clean_line = clean_line.strip("#*` ")
        if clean_line:
            bullets.append(clean_line)
        if len(bullets) >= max_bullets:
            break

    if bullets:
        return [_clean_text(bullet, 105) for bullet in bullets]

    sentences = re.split(r"(?<=[.!?])\s+", content.strip())
    return [
        _clean_text(sentence, 105)
        for sentence in sentences
        if sentence
    ][:max_bullets]


def _bullets_to_narration(
    heading: str,
    bullets: list[str],
    video_instructions: str | None,
) -> str:
    narration = f"{heading}. " + " ".join(bullets)
    if video_instructions and video_instructions.strip():
        narration += f" Presentation note: {video_instructions.strip()[:500]}"
    return narration


def _build_timestamped_script(slides: list[VideoSlide]) -> list[TimestampedScriptLine]:
    script = []
    cursor = 0
    for slide in slides:
        end = cursor + slide.duration_seconds
        script.append(
            TimestampedScriptLine(
                start_seconds=cursor,
                end_seconds=end,
                slide_title=slide.title,
                narration=slide.narration,
                subtitle=slide.subtitle,
            )
        )
        cursor = end
    return script


def _apply_voiceover_durations(
    slides: list[VideoSlide],
    durations: list[int],
) -> list[VideoSlide]:
    updated_slides = []
    for index, slide in enumerate(slides):
        if index < len(durations):
            updated_slides.append(
                replace(slide, duration_seconds=max(2, durations[index]))
            )
        else:
            updated_slides.append(slide)
    return updated_slides


def _write_presentation(slides: list[VideoSlide], pptx_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise VideoSummaryError("python-pptx is not installed.") from exc

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    for slide_number, slide_spec in enumerate(slides, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            presentation.slide_width,
            presentation.slide_height,
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(249, 247, 242)
        background.line.fill.background()

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(0.18),
            presentation.slide_height,
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(13, 116, 109)
        accent.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.42), Inches(6.45), Inches(0.78))
        title_frame = title_box.text_frame
        title_frame.text = slide_spec.title
        title_frame.paragraphs[0].font.size = Pt(27)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)

        body_box = slide.shapes.add_textbox(Inches(0.68), Inches(1.62), Inches(5.95), Inches(4.7))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.margin_left = 0
        body_frame.margin_right = 0
        for index, bullet in enumerate(slide_spec.bullets[:3]):
            paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = RGBColor(31, 41, 55)
            paragraph.space_after = Pt(14)

        visual_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.05),
            Inches(1.18),
            Inches(5.65),
            Inches(4.72),
        )
        visual_box.fill.solid()
        visual_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        visual_box.line.color.rgb = RGBColor(226, 232, 240)

        if slide_spec.visual_path and Path(slide_spec.visual_path).exists():
            _add_fit_picture(slide, slide_spec.visual_path, Inches(7.25), Inches(1.36), Inches(5.25), Inches(4.18))
        else:
            placeholder = slide.shapes.add_textbox(Inches(7.45), Inches(2.65), Inches(4.8), Inches(0.95))
            placeholder.text_frame.text = "Conceptual overview"
            placeholder.text_frame.paragraphs[0].font.size = Pt(24)
            placeholder.text_frame.paragraphs[0].font.bold = True
            placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(13, 116, 109)

        if slide_spec.visual_caption:
            caption_box = slide.shapes.add_textbox(Inches(7.15), Inches(6.02), Inches(5.45), Inches(0.48))
            caption_box.text_frame.text = slide_spec.visual_caption
            caption_box.text_frame.paragraphs[0].font.size = Pt(10.5)
            caption_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)

        footer = slide.shapes.add_textbox(Inches(0.65), Inches(6.86), Inches(6.5), Inches(0.28))
        footer.text_frame.text = f"Slide {slide_number} / {len(slides)}"
        footer.text_frame.paragraphs[0].font.size = Pt(9.5)
        footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)

    presentation.save(pptx_path)


def _add_fit_picture(slide, image_path: str, left, top, max_width, max_height) -> None:
    try:
        from PIL import Image
    except ImportError as exc:
        raise VideoSummaryError("Pillow is not installed.") from exc

    try:
        with Image.open(image_path) as image:
            image_width, image_height = image.size
    except OSError:
        return

    ratio = min(max_width / image_width, max_height / image_height)
    width = int(image_width * ratio)
    height = int(image_height * ratio)
    centered_left = left + int((max_width - width) / 2)
    centered_top = top + int((max_height - height) / 2)
    slide.shapes.add_picture(image_path, centered_left, centered_top, width=width, height=height)


def _render_slide_images(slides: list[VideoSlide], frames_dir: Path) -> list[Path]:
    try:
        from PIL import Image, ImageDraw, ImageFont, ImageOps
    except ImportError as exc:
        raise VideoSummaryError("Pillow is not installed.") from exc

    title_font = _load_font(ImageFont, 42)
    bullet_font = _load_font(ImageFont, 27)
    small_font = _load_font(ImageFont, 20)
    caption_font = _load_font(ImageFont, 16)
    frame_paths = []

    for index, slide in enumerate(slides, start=1):
        image = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), "#f9f7f2")
        draw = ImageDraw.Draw(image)
        draw.rectangle((0, 0, 20, SLIDE_HEIGHT), fill="#0d746d")
        draw.rectangle((60, 55, 90, 61), fill="#d97706")
        _draw_wrapped_text(draw, slide.title, (60, 78), title_font, "#1f2937", 31, 2, line_gap=9)

        y = 190
        for bullet in slide.bullets[:3]:
            draw.ellipse((65, y + 10, 76, y + 21), fill="#0d746d")
            consumed = _draw_wrapped_text(
                draw,
                bullet,
                (96, y),
                bullet_font,
                "#1f2937",
                34,
                3,
                line_gap=8,
            )
            y += consumed + 28

        visual_box = (720, 118, 1228, 575)
        draw.rounded_rectangle((710, 108, 1238, 585), radius=8, fill="#e5e7eb")
        draw.rounded_rectangle(visual_box, radius=8, fill="#ffffff", outline="#d1d5db", width=2)

        if slide.visual_path and Path(slide.visual_path).exists():
            try:
                visual = Image.open(slide.visual_path).convert("RGB")
                visual.thumbnail((visual_box[2] - visual_box[0] - 34, visual_box[3] - visual_box[1] - 34))
                paste_x = visual_box[0] + ((visual_box[2] - visual_box[0]) - visual.width) // 2
                paste_y = visual_box[1] + ((visual_box[3] - visual_box[1]) - visual.height) // 2
                image.paste(visual, (paste_x, paste_y))
            except OSError:
                _draw_placeholder_visual(draw, visual_box, small_font)
        else:
            _draw_placeholder_visual(draw, visual_box, small_font)

        if slide.visual_caption:
            _draw_wrapped_text(
                draw,
                slide.visual_caption,
                (725, 604),
                caption_font,
                "#64748b",
                58,
                2,
                line_gap=4,
            )

        footer = f"Slide {index} of {len(slides)}  |  {slide.duration_seconds}s"
        draw.text((60, SLIDE_HEIGHT - 48), footer, fill="#64748b", font=small_font)

        frame_path = frames_dir / f"slide_{index:02d}.png"
        image.save(frame_path)
        frame_paths.append(frame_path)

    return frame_paths


def _draw_placeholder_visual(draw, visual_box: tuple[int, int, int, int], font) -> None:
    x1, y1, x2, y2 = visual_box
    draw.line((x1 + 70, y2 - 90, x1 + 180, y1 + 185, x1 + 290, y2 - 145, x2 - 70, y1 + 130), fill="#0d746d", width=5)
    draw.rectangle((x1 + 70, y1 + 110, x2 - 70, y2 - 70), outline="#94a3b8", width=3)
    draw.text((x1 + 135, y2 - 58), "Visual generated from paper structure", fill="#64748b", font=font)


def _draw_wrapped_text(
    draw,
    text: str,
    position: tuple[int, int],
    font,
    fill: str,
    wrap_width: int,
    max_lines: int,
    *,
    line_gap: int,
) -> int:
    x, y = position
    line_height = _font_height(font) + line_gap
    lines = wrap(text, width=wrap_width)[:max_lines]
    for line in lines:
        draw.text((x, y), line, fill=fill, font=font)
        y += line_height
    return max(1, len(lines)) * line_height


def _write_script(script: list[TimestampedScriptLine], script_path: Path) -> None:
    lines = []
    for line in script:
        lines.append(
            f"[{_format_short_timestamp(line.start_seconds)} - "
            f"{_format_short_timestamp(line.end_seconds)}] "
            f"{line.slide_title}: {line.narration}"
        )
    script_path.write_text("\n\n".join(lines), encoding="utf-8")


def _write_script_json(script: list[TimestampedScriptLine], script_json_path: Path) -> None:
    script_json_path.write_text(
        json.dumps([asdict(line) for line in script], indent=2),
        encoding="utf-8",
    )


def _write_vtt(script: list[TimestampedScriptLine], subtitle_vtt_path: Path) -> None:
    blocks = ["WEBVTT"]
    for line in script:
        blocks.append(
            "\n".join(
                [
                    f"{_format_vtt_timestamp(line.start_seconds)} --> {_format_vtt_timestamp(line.end_seconds)}",
                    line.subtitle,
                ]
            )
        )
    subtitle_vtt_path.write_text("\n\n".join(blocks) + "\n", encoding="utf-8")


def _write_srt(script: list[TimestampedScriptLine], subtitle_srt_path: Path) -> None:
    blocks = []
    for index, line in enumerate(script, start=1):
        blocks.append(
            "\n".join(
                [
                    str(index),
                    f"{_format_srt_timestamp(line.start_seconds)} --> {_format_srt_timestamp(line.end_seconds)}",
                    line.subtitle,
                ]
            )
        )
    subtitle_srt_path.write_text("\n\n".join(blocks) + "\n", encoding="utf-8")


def _write_asset_metadata(image_assets: list, asset_metadata_path: Path) -> None:
    asset_metadata_path.write_text(
        json.dumps([asdict(asset) for asset in image_assets], indent=2),
        encoding="utf-8",
    )


def _render_video(
    slide_image_paths: list[Path],
    slides: list[VideoSlide],
    video_path: Path,
) -> None:
    try:
        import imageio.v2 as imageio
    except ImportError as exc:
        raise VideoSummaryError("imageio is not installed.") from exc

    with imageio.get_writer(video_path, fps=VIDEO_FPS, codec="libx264") as writer:
        for image_path, slide in zip(slide_image_paths, slides):
            frame = imageio.imread(image_path)
            for _ in range(max(1, slide.duration_seconds * VIDEO_FPS)):
                writer.append_data(frame)


def _mux_video_audio(
    video_path: Path,
    audio_path: Path,
    output_path: Path,
) -> None:
    try:
        import imageio_ffmpeg
    except ImportError as exc:
        raise VideoSummaryError("imageio-ffmpeg is not installed.") from exc

    command = [
        imageio_ffmpeg.get_ffmpeg_exe(),
        "-y",
        "-i",
        str(video_path),
        "-i",
        str(audio_path),
        "-c:v",
        "copy",
        "-c:a",
        "aac",
        "-b:a",
        "128k",
        str(output_path),
    ]
    try:
        subprocess.run(command, check=True, capture_output=True, text=True)
    except (OSError, subprocess.CalledProcessError) as exc:
        raise VideoSummaryError("Adding voiceover audio to the video failed.") from exc


def _estimate_duration(narration: str, minimum_seconds: int) -> int:
    word_count = len(narration.split())
    return min(max(minimum_seconds, round(word_count / 2.4)), 24)


def _coerce_duration(value, default_seconds: int) -> int:
    try:
        return _clamp_duration(int(value))
    except (TypeError, ValueError):
        return default_seconds


def _clamp_duration(seconds: int) -> int:
    return min(max(seconds, 5), 20)


def _coerce_int(value) -> int | None:
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _format_short_timestamp(seconds: int) -> str:
    minutes, remaining_seconds = divmod(seconds, 60)
    return f"{minutes:02d}:{remaining_seconds:02d}"


def _format_vtt_timestamp(seconds: int) -> str:
    hours, remainder = divmod(seconds, 3600)
    minutes, remaining_seconds = divmod(remainder, 60)
    return f"{hours:02d}:{minutes:02d}:{remaining_seconds:02d}.000"


def _format_srt_timestamp(seconds: int) -> str:
    hours, remainder = divmod(seconds, 3600)
    minutes, remaining_seconds = divmod(remainder, 60)
    return f"{hours:02d}:{minutes:02d}:{remaining_seconds:02d},000"


def _clean_bullets(value) -> list[str]:
    if not isinstance(value, list):
        return []
    bullets = []
    for bullet in value[:3]:
        clean_bullet = _clean_text(bullet, limit=105)
        if clean_bullet:
            bullets.append(clean_bullet)
    return bullets


def _clean_text(value, limit: int) -> str:
    if value is None:
        return ""
    clean_text = " ".join(str(value).split())
    if len(clean_text) <= limit:
        return clean_text
    return clean_text[: limit - 3].rstrip() + "..."


def _short_caption(text: str, *, limit: int = 220) -> str:
    first_sentence = re.split(r"(?<=[.!?])\s+", text.strip())[0]
    return _clean_text(first_sentence, limit=limit)


def _load_font(image_font_module, size: int):
    for font_name in ("Arial.ttf", "Helvetica.ttc", "DejaVuSans.ttf"):
        try:
            return image_font_module.truetype(font_name, size)
        except OSError:
            continue
    return image_font_module.load_default()


def _font_height(font) -> int:
    try:
        return int(font.getbbox("Ag")[3] - font.getbbox("Ag")[1])
    except AttributeError:
        return int(font.getsize("Ag")[1])

# [GenAI Usage] Reflection: sometimes the openAI models itself is not able to write good slides by itself. 
# extra plans are necessary as intermediate procedures. 